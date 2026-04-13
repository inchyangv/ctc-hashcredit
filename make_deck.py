#!/usr/bin/env python3
"""
HashCredit Pitch Deck generator — python-pptx
Style: Swiss International (07)
  - White background, near-black text
  - Single accent: Bitcoin Orange #F7931A
  - Left vertical accent bar, horizontal dividers, strict grid
  - Restrained typography: Arial Bold titles, Arial body
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TOTAL_SLIDES = 15

# ── Swiss International palette ──────────────────────────────────────────────
BG       = RGBColor(0xFF, 0xFF, 0xFF)   # pure white
BLACK    = RGBColor(0x11, 0x11, 0x11)   # near-black
DARK     = RGBColor(0x33, 0x33, 0x33)   # headings
BODY     = RGBColor(0x55, 0x55, 0x55)   # body text
MUTED    = RGBColor(0x99, 0x99, 0x99)   # captions
ACCENT   = RGBColor(0xF7, 0x93, 0x1A)   # Bitcoin orange — the ONE accent
NAVY     = RGBColor(0x1A, 0x1A, 0x2E)   # dark card background
LTGRAY   = RGBColor(0xF5, 0xF5, 0xF5)   # light card fill
RULE     = RGBColor(0xDD, 0xDD, 0xDD)   # horizontal rules
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED_SOFT = RGBColor(0xCC, 0x33, 0x33)   # negative items
GREEN_OK = RGBColor(0x22, 0x88, 0x55)   # positive items

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Layout constants (strict grid) ───────────────────────────────────────────
MARGIN_L = Inches(0.9)
MARGIN_R = Inches(0.9)
MARGIN_T = Inches(0.7)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
BAR_W    = Inches(0.08)          # left vertical accent bar

# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def tb(slide, l, t, w, h, text, size=Pt(14), bold=False,
       color=BODY, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Arial"
    return txb

def add_para(tf, text, size=Pt(13), bold=False, color=BODY,
             align=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return p

# ── Signature elements ───────────────────────────────────────────────────────

def left_bar(slide, color=ACCENT):
    """Swiss International signature: vertical left accent bar."""
    add_rect(slide, Inches(0.45), Inches(0.5),
             BAR_W, SLIDE_H - Inches(1.0), fill_color=color)

def h_rule(slide, y):
    """Horizontal divider rule."""
    add_rect(slide, MARGIN_L, y, CONTENT_W, Pt(1), fill_color=RULE)

def slide_num(slide, num):
    tb(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
       Inches(1.0), Inches(0.3),
       f"{num}", size=Pt(10), color=MUTED, align=PP_ALIGN.RIGHT)

def section_label(slide, text):
    """Top-left section label (small caps feel)."""
    tb(slide, MARGIN_L, Inches(0.55), Inches(4), Inches(0.3),
       text.upper(), size=Pt(9), bold=True, color=ACCENT)

def heading(slide, text, y=Inches(1.0), size=Pt(32)):
    """Main slide heading."""
    tb(slide, MARGIN_L, y, CONTENT_W, Inches(0.7),
       text, size=size, bold=True, color=BLACK)

def subhead(slide, text, y=Inches(1.6)):
    """Subtitle / description line."""
    tb(slide, MARGIN_L, y, Inches(9.5), Inches(0.5),
       text, size=Pt(14), color=BODY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═════════════════════════════════════════════════════════════════════════════

# ── 1. COVER ─────────────────────────────────────────────────────────────────
def slide_cover(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    # Accent bar — full height, thicker for cover
    add_rect(sl, Inches(0.45), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT)

    # Title
    tb(sl, Inches(1.2), Inches(2.0), Inches(10), Inches(1.0),
       "HashCredit", size=Pt(72), bold=True, color=BLACK)

    # Tagline
    tb(sl, Inches(1.2), Inches(3.2), Inches(10), Inches(0.6),
       "Working capital for Bitcoin miners, via mining pools.",
       size=Pt(24), color=DARK)

    # Three descriptors
    h_rule(sl, Inches(4.2))
    descs = [
        "Pool-enforced repayment",
        "SPV-proven revenue",
        "USDT on demand",
    ]
    for i, d in enumerate(descs):
        x = Inches(1.2) + i * Inches(3.8)
        tb(sl, x, Inches(4.5), Inches(3.5), Inches(0.35),
           d, size=Pt(13), bold=True, color=BODY)

    # Bottom bar
    h_rule(sl, Inches(5.8))
    tb(sl, Inches(1.2), Inches(6.0), Inches(10), Inches(0.3),
       "Built on HashKey Chain  ·  DeFi Track  ·  Live on Testnet",
       size=Pt(11), color=MUTED)

    slide_num(sl, 1)


# ── 2. PROBLEM ───────────────────────────────────────────────────────────────
def slide_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "01 — Problem")
    heading(sl, "$17B in annual mining revenue.")
    tb(sl, MARGIN_L, Inches(1.55), Inches(9), Inches(0.4),
       "Zero ways to borrow against it on-chain.",
       size=Pt(22), bold=True, color=ACCENT)

    subhead(sl, "Miners pay for electricity, hardware, and facilities in fiat — but earn in BTC.",
            y=Inches(2.1))

    h_rule(sl, Inches(2.7))

    # Pain points — clean list
    pain = [
        ("Lock up your BTC?",       "Ledn / SALT — 7–12% APR, requires BTC collateral lockup."),
        ("Pledge your ASICs?",      "Luxor Finance — over-collateralized, institutional-only."),
        ("Raise equity?",           "Dilutive, takes months."),
        ("Revenue-based off-chain?","No pool enforcement, no on-chain verifiability."),
    ]
    for i, (h, d) in enumerate(pain):
        y = Inches(3.0) + i * Inches(0.75)
        tb(sl, MARGIN_L, y, Inches(0.2), Inches(0.3),
           "×", size=Pt(14), bold=True, color=RED_SOFT)
        tb(sl, Inches(1.2), y, Inches(2.8), Inches(0.3),
           h, size=Pt(13), bold=True, color=BLACK)
        tb(sl, Inches(4.2), y, Inches(5.0), Inches(0.3),
           d, size=Pt(12), color=BODY)

    # Right side — key stats
    for i, (val, lbl) in enumerate([
        ("$17.2B", "Annual BTC miner revenue"),
        ("$11B+",  "Mining debt since 2023"),
        ("1,200+", "Days avg hardware ROI"),
    ]):
        y = Inches(1.2) + i * Inches(1.7)
        add_rect(sl, Inches(10.0), y, Inches(2.8), Inches(1.3), fill_color=LTGRAY)
        tb(sl, Inches(10.2), y + Inches(0.15), Inches(2.4), Inches(0.6),
           val, size=Pt(32), bold=True, color=ACCENT)
        tb(sl, Inches(10.2), y + Inches(0.75), Inches(2.4), Inches(0.4),
           lbl, size=Pt(10), color=MUTED)

    slide_num(sl, 2)


# ── 3. MARKET ────────────────────────────────────────────────────────────────
def slide_market(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "02 — Market Opportunity")
    heading(sl, "$17.2B TAM — targeting where the need is greatest.")
    subhead(sl, "Mid-market miners: 30–40% of global hashrate. Too small for capital markets, too big for personal credit.")

    h_rule(sl, Inches(2.3))

    # TAM → SAM → SOM
    for i, (val, title, desc) in enumerate([
        ("$17.2B", "TAM", "Global BTC miner annual revenue (2025)"),
        ("$1.7B",  "SAM", "10% market penetration, mid-market focus"),
        ("$136M",  "SOM", "Protocol revenue at 8% APR"),
    ]):
        x = MARGIN_L + i * Inches(4.0)
        add_rect(sl, x, Inches(2.6), Inches(3.6), Inches(2.0), fill_color=LTGRAY)
        tb(sl, x + Inches(0.3), Inches(2.8), Inches(3.0), Inches(0.6),
           val, size=Pt(36), bold=True, color=ACCENT if i == 0 else BLACK)
        tb(sl, x + Inches(0.3), Inches(3.4), Inches(3.0), Inches(0.3),
           title, size=Pt(11), bold=True, color=DARK)
        tb(sl, x + Inches(0.3), Inches(3.75), Inches(3.0), Inches(0.5),
           desc, size=Pt(11), color=BODY)

    # LP opportunity
    h_rule(sl, Inches(5.0))
    tb(sl, MARGIN_L, Inches(5.3), Inches(11), Inches(0.4),
       "LP opportunity: USDT depositors earn 8% APR — 2–3× standard DeFi rates (Aave USDT ~3–4%), backed by SPV-proven revenue.",
       size=Pt(13), color=DARK)

    # Source
    tb(sl, MARGIN_L, Inches(6.6), Inches(6), Inches(0.3),
       "Sources: The Block, Precedence Research, DeFi Llama", size=Pt(9), color=MUTED)

    slide_num(sl, 3)


# ── 4. INSIGHT ───────────────────────────────────────────────────────────────
def slide_insight(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "03 — Key Insight")

    # Big quote-style statement
    tb(sl, MARGIN_L, Inches(1.4), Inches(10.5), Inches(1.2),
       "You can't prove hashrate directly —\nit's a physical rate.\nBut you can prove its output.",
       size=Pt(36), bold=True, color=BLACK)

    h_rule(sl, Inches(3.3))

    tb(sl, MARGIN_L, Inches(3.6), Inches(10.5), Inches(1.2),
       "Every pool payout is a Bitcoin transaction proportional to contributed hash power. "
       "It's timestamped, immutable, and verifiable by anyone with block headers.\n\n"
       "Payout history is the hashrate record — and SPV turns it into trustless on-chain evidence.",
       size=Pt(14), color=BODY)

    # Comparison box
    add_rect(sl, MARGIN_L, Inches(5.4), Inches(11.5), Inches(1.2), fill_color=LTGRAY)
    tb(sl, Inches(1.2), Inches(5.55), Inches(5), Inches(0.3),
       "Stripe Capital", size=Pt(12), bold=True, color=BODY)
    tb(sl, Inches(1.2), Inches(5.85), Inches(5), Inches(0.5),
       "Self-reported revenue through their platform.\nTrust-based.", size=Pt(11), color=MUTED)

    tb(sl, Inches(6.8), Inches(5.55), Inches(5), Inches(0.3),
       "HashCredit", size=Pt(12), bold=True, color=ACCENT)
    tb(sl, Inches(6.8), Inches(5.85), Inches(5), Inches(0.5),
       "Bitcoin payout verified by proof-of-work.\nMath-based.", size=Pt(11), color=DARK)

    slide_num(sl, 4)


# ── 5. SOLUTION ──────────────────────────────────────────────────────────────
def slide_solution(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "04 — Solution")
    heading(sl, "Verified Bitcoin payouts → revolving USDT credit line.")

    # Flow
    tb(sl, MARGIN_L, Inches(1.7), Inches(11), Inches(0.5),
       "Pool payout  →  SPV proof  →  Credit limit update  →  Draw USDT  →  Pool auto-withholds repayment",
       size=Pt(14), bold=True, color=DARK)

    h_rule(sl, Inches(2.4))

    bullets = [
        ("No BTC lockup", "Mine through your registered pool — no collateral required."),
        ("Real-time credit", "Credit limit adjusts automatically with each verified payout."),
        ("Auto-repayment", "Pool withholds a percentage of each subsequent payout."),
        ("Fully on-chain", "Every step auditable. No intermediary."),
        ("Modular architecture", "IVerifierAdapter — swap proof sources without touching credit logic."),
    ]
    for i, (h, d) in enumerate(bullets):
        y = Inches(2.7) + i * Inches(0.72)
        # Orange bullet marker
        add_rect(sl, MARGIN_L, y + Inches(0.05), Inches(0.06), Inches(0.06), fill_color=ACCENT)
        tb(sl, Inches(1.2), y, Inches(3.0), Inches(0.3),
           h, size=Pt(13), bold=True, color=BLACK)
        tb(sl, Inches(4.5), y, Inches(7.0), Inches(0.3),
           d, size=Pt(12), color=BODY)

    # One-liner box
    add_rect(sl, MARGIN_L, Inches(6.3), CONTENT_W, Inches(0.6), fill_color=LTGRAY)
    tb(sl, Inches(1.2), Inches(6.38), Inches(10), Inches(0.4),
       '"Stripe Capital for Bitcoin miners — but trustless, and pool-enforced."',
       size=Pt(13), italic=True, color=DARK, align=PP_ALIGN.CENTER)

    slide_num(sl, 5)


# ── 6. HOW IT WORKS ──────────────────────────────────────────────────────────
def slide_how(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "05 — How It Works")
    heading(sl, "Seven steps. Fully automated.")

    steps = [
        ("1", "Pool registers", "Agrees to withhold repayment from payouts."),
        ("2", "Miner gets paid", "Bitcoin transaction from pool to miner address."),
        ("3", "Worker detects", "Off-chain worker monitors BTC addresses."),
        ("4", "SPV proof built", "Headers + Merkle inclusion + tx output."),
        ("5", "On-chain verify", "PoW check + Merkle check + output script match."),
        ("6", "Credit updates", "Trailing-window credit limit recalculated."),
        ("7", "Draw / Repay", "Miner draws USDT. Pool withholds from next payout."),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(1.7) + i * Inches(0.73)
        # Step number
        add_rect(sl, MARGIN_L, y, Inches(0.5), Inches(0.5), fill_color=ACCENT)
        tb(sl, MARGIN_L, y + Inches(0.05), Inches(0.5), Inches(0.4),
           num, size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title + desc
        tb(sl, Inches(1.6), y + Inches(0.05), Inches(2.5), Inches(0.35),
           title, size=Pt(13), bold=True, color=BLACK)
        tb(sl, Inches(4.3), y + Inches(0.05), Inches(6.0), Inches(0.35),
           desc, size=Pt(12), color=BODY)

    # Enforcement note
    h_rule(sl, Inches(6.8))
    tb(sl, MARGIN_L, Inches(6.9), Inches(11), Inches(0.3),
       "Default enforcement: pool redirects miner's hashrate. No courts. Pure economic enforcement.",
       size=Pt(11), bold=True, color=DARK)

    slide_num(sl, 6)


# ── 7. BUSINESS MODEL ────────────────────────────────────────────────────────
def slide_biz_model(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "06 — Business Model")
    heading(sl, "Three revenue streams.")

    h_rule(sl, Inches(1.6))

    models = [
        ("Interest spread", "Borrowers pay 10% APR → LPs earn 8% → protocol keeps 2% spread."),
        ("Origination fee", "0.5% per drawdown."),
        ("Pool withholding", "Automatic repayment via pool — zero collection cost."),
    ]
    for i, (title, desc) in enumerate(models):
        y = Inches(1.9) + i * Inches(0.8)
        add_rect(sl, MARGIN_L, y + Inches(0.05), Inches(0.06), Inches(0.06), fill_color=ACCENT)
        tb(sl, Inches(1.2), y, Inches(3.0), Inches(0.3),
           title, size=Pt(14), bold=True, color=BLACK)
        tb(sl, Inches(4.5), y, Inches(7.0), Inches(0.3),
           desc, size=Pt(12), color=BODY)

    h_rule(sl, Inches(4.4))

    # LP yield comparison
    tb(sl, MARGIN_L, Inches(4.7), Inches(4), Inches(0.3),
       "LP Yield Comparison", size=Pt(14), bold=True, color=BLACK)

    rows = [
        ("Platform",      "APR",     "Security"),
        ("Aave USDT",     "3–4%",    "Overcollateral"),
        ("Curve stables", "5–7%",    "Overcollateral + slippage"),
        ("HashCredit LP", "8%",      "SPV-proven revenue + pool withholding"),
    ]
    for i, (a, b, c) in enumerate(rows):
        y = Inches(5.1) + i * Inches(0.42)
        is_header = (i == 0)
        is_ours = (i == 3)
        col = MUTED if is_header else (ACCENT if is_ours else BODY)
        w = Pt(11) if is_header else Pt(12)
        tb(sl, MARGIN_L, y, Inches(2.5), Inches(0.3),
           a, size=w, bold=is_header or is_ours, color=col)
        tb(sl, Inches(3.6), y, Inches(1.5), Inches(0.3),
           b, size=w, bold=is_header or is_ours, color=col)
        tb(sl, Inches(5.5), y, Inches(5.0), Inches(0.3),
           c, size=w, bold=is_header, color=col)

    slide_num(sl, 7)


# ── 8. TRACTION ──────────────────────────────────────────────────────────────
def slide_traction(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "07 — Traction")
    heading(sl, "It works today. End-to-end on testnet.")

    h_rule(sl, Inches(1.6))

    items = [
        "7 production smart contracts live on HashKey Chain testnet",
        "SPV proofs generated from real Bitcoin testnet transactions",
        "Full borrow/repay lifecycle operational",
        "Automated prover worker running 24/7",
        "Frontend dashboard live",
        "Test suite: unit, integration, invariant fuzzing, gas profiling",
        "Modular IVerifierAdapter — proof source is pluggable",
    ]
    for i, item in enumerate(items):
        y = Inches(1.9) + i * Inches(0.6)
        tb(sl, MARGIN_L, y, Inches(0.3), Inches(0.3),
           "✓", size=Pt(14), bold=True, color=GREEN_OK)
        tb(sl, Inches(1.3), y, Inches(10), Inches(0.3),
           item, size=Pt(13), color=BLACK)

    # Contract info
    h_rule(sl, Inches(6.3))
    tb(sl, MARGIN_L, Inches(6.45), Inches(10), Inches(0.3),
       "HashKey Chain Testnet · chainId 133 · Verifiable now",
       size=Pt(11), color=MUTED)

    slide_num(sl, 8)


# ── 9. WHY NOW ───────────────────────────────────────────────────────────────
def slide_why_now(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "08 — Why Now")
    heading(sl, "Structural demand. No incumbent.")

    h_rule(sl, Inches(1.6))

    signals = [
        ("Halving (2024)",     "Block reward → 3.125 BTC. Revenue per hash is down. Working capital matters more."),
        ("Hashrate at 1 ZH/s", "Hardware ROI now exceeds 1,200 days. Miners can't afford to wait."),
        ("Treasury drawdowns",  "Public miners sold 5,359 BTC in Dec 2025 alone to stay liquid."),
        ("No incumbent",        "Trustless BTC proof + pool-enforced repayment — nobody else is doing this."),
    ]
    for i, (signal, desc) in enumerate(signals):
        y = Inches(1.9) + i * Inches(1.1)
        add_rect(sl, MARGIN_L, y, CONTENT_W, Inches(0.85), fill_color=LTGRAY)
        tb(sl, Inches(1.2), y + Inches(0.1), Inches(3.0), Inches(0.3),
           signal, size=Pt(14), bold=True, color=BLACK)
        tb(sl, Inches(1.2), y + Inches(0.45), Inches(10.5), Inches(0.3),
           desc, size=Pt(12), color=BODY)

    slide_num(sl, 9)


# ── 10. COMPETITIVE EDGE ────────────────────────────────────────────────────
def slide_competitive(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "09 — Competitive Edge")
    heading(sl, "Cryptographic proof, not trust. Pool-enforced, not manual.")

    h_rule(sl, Inches(1.6))

    headers = ["", "HashCredit", "Everyone else"]
    rows = [
        ("Collateral",    "None — revenue-proven",          "BTC, ASICs, or financials"),
        ("Verification",  "SPV proof — pure math",          "Oracle, custodian, or manual"),
        ("Repayment",     "Automatic — pool withholds",     "Manual, trust-dependent"),
        ("Speed",         "Instant — auto-updating credit", "Days to weeks"),
        ("Access",        "Permissionless via pool",        "KYC, minimums, gatekeeping"),
    ]

    # Header row
    y = Inches(1.9)
    for j, h in enumerate(headers):
        x = MARGIN_L + j * Inches(4.0)
        tb(sl, x, y, Inches(3.8), Inches(0.3),
           h, size=Pt(11), bold=True, color=MUTED)

    # Data rows
    for i, (label, ours, theirs) in enumerate(rows):
        y = Inches(2.4) + i * Inches(0.8)
        if i % 2 == 0:
            add_rect(sl, MARGIN_L, y - Inches(0.05), CONTENT_W, Inches(0.7), fill_color=LTGRAY)
        tb(sl, MARGIN_L, y, Inches(3.8), Inches(0.5),
           label, size=Pt(12), bold=True, color=BLACK)
        tb(sl, MARGIN_L + Inches(4.0), y, Inches(3.8), Inches(0.5),
           ours, size=Pt(12), bold=True, color=ACCENT)
        tb(sl, MARGIN_L + Inches(8.0), y, Inches(3.8), Inches(0.5),
           theirs, size=Pt(12), color=BODY)

    h_rule(sl, Inches(6.5))
    tb(sl, MARGIN_L, Inches(6.65), Inches(11), Inches(0.3),
       "Same trust model Bitcoin itself has used since 2009. No oracle to bribe. No custodian to trust.",
       size=Pt(11), color=DARK)

    slide_num(sl, 10)


# ── 11. ROADMAP ──────────────────────────────────────────────────────────────
def slide_roadmap(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "10 — Roadmap")
    heading(sl, "From testnet to $5M TVL in 12 months.")

    h_rule(sl, Inches(1.6))

    phases = [
        ("Q1 2026", "Complete",    "7 contracts, SPV verifier, frontend, full lifecycle on testnet.", True),
        ("Q2 2026", "In Progress", "Security audit, HashKey Chain mainnet, 10 pilot miners.", False),
        ("Q3 2026", "Target",      "50 miners, $500K TVL, mining pool API partnerships.", False),
        ("Q4 2026", "Vision",      "$5M TVL, 200+ miners, cross-chain oracle adapter.", False),
    ]
    for i, (when, status, desc, done) in enumerate(phases):
        y = Inches(2.0) + i * Inches(1.15)
        # Timeline dot
        dot_color = GREEN_OK if done else ACCENT
        add_rect(sl, Inches(1.2), y + Inches(0.12), Inches(0.14), Inches(0.14), fill_color=dot_color)
        # Vertical connector
        if i < len(phases) - 1:
            add_rect(sl, Inches(1.26), y + Inches(0.3), Inches(0.02), Inches(0.85), fill_color=RULE)

        tb(sl, Inches(1.6), y, Inches(1.5), Inches(0.3),
           when, size=Pt(13), bold=True, color=BLACK)
        tb(sl, Inches(3.2), y, Inches(1.5), Inches(0.3),
           status, size=Pt(11), bold=True, color=GREEN_OK if done else ACCENT)
        tb(sl, Inches(1.6), y + Inches(0.35), Inches(9.5), Inches(0.5),
           desc, size=Pt(12), color=BODY)

    slide_num(sl, 11)


# ── 12. TEAM ─────────────────────────────────────────────────────────────────
def slide_team(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "11 — Team")
    heading(sl, "Built by engineers who understand both Bitcoin and DeFi.")

    h_rule(sl, Inches(1.6))

    team = [
        ("Incheol Yang", "Co-Founder",
         "Co-founded a DeFi system trading house, managed $20M AUM at 40%+ APR.\n"
         "Previously at KRAFTON (PUBG) and Coinone Exchange building trading infrastructure.\n"
         "Full-stack protocol design: SPV verifier, credit engine, vault, frontend."),
        ("Juhyeong Park", "Co-Founder",
         "CTO of Onther — led a mainnet to $750M market cap.\n"
         "Designed Plasma EVM. Smart contract audits and full-stack Solidity architecture.\n"
         "Core contract engineering: HashCreditManager, LendingVault, RiskConfig."),
    ]
    for i, (name, role, bio) in enumerate(team):
        y = Inches(2.0) + i * Inches(2.5)
        add_rect(sl, MARGIN_L, y, Inches(11.5), Inches(2.1), fill_color=LTGRAY)

        tb(sl, Inches(1.2), y + Inches(0.2), Inches(4), Inches(0.4),
           name, size=Pt(20), bold=True, color=BLACK)
        tb(sl, Inches(1.2), y + Inches(0.65), Inches(2), Inches(0.3),
           role, size=Pt(12), bold=True, color=ACCENT)
        tb(sl, Inches(4.0), y + Inches(0.2), Inches(8.0), Inches(1.7),
           bio, size=Pt(12), color=BODY)

    slide_num(sl, 12)


# ── 13. WHY HASHKEY CHAIN ────────────────────────────────────────────────────
def slide_why_hashkey(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "12 — Why HashKey Chain")
    heading(sl, "Mining-revenue lending needs more than an EVM chain.")

    h_rule(sl, Inches(1.6))

    points = [
        ("Compliance infrastructure",
         "Pool withholding requires contractual + legal framework. "
         "HashKey Group (SFC-licensed HK) provides identity tooling, KYC, and legal enforceability."),
        ("Ecosystem as distribution",
         "HashKey Exchange → fiat on/off-ramp + USDT liquidity. "
         "HashKey Capital → mining operator introductions. "
         "Hong Kong → center of Asia-Pacific hashrate (~30% of global)."),
        ("Full EVM precompile support",
         "ecrecover + sha256 + ripemd160 enable trustless on-chain BTC address verification. "
         "HashKey Chain (OP Stack) supports all standard precompiles natively. No oracle. No bridge."),
    ]
    for i, (title, desc) in enumerate(points):
        y = Inches(1.9) + i * Inches(1.55)
        add_rect(sl, MARGIN_L, y, CONTENT_W, Inches(1.25), fill_color=LTGRAY)
        # Orange number
        tb(sl, Inches(1.1), y + Inches(0.12), Inches(0.5), Inches(0.4),
           str(i + 1), size=Pt(24), bold=True, color=ACCENT)
        tb(sl, Inches(1.7), y + Inches(0.15), Inches(4.0), Inches(0.3),
           title, size=Pt(14), bold=True, color=BLACK)
        tb(sl, Inches(1.7), y + Inches(0.5), Inches(10.0), Inches(0.6),
           desc, size=Pt(12), color=BODY)

    slide_num(sl, 13)


# ── 14. THE ASK ──────────────────────────────────────────────────────────────
def slide_ask(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    left_bar(sl)
    section_label(sl, "13 — The Ask")
    heading(sl, "Raising $250K Seed.")

    subhead(sl, "Revenue-based financing for Bitcoin miners — trustless, pool-enforced, production-ready.")

    h_rule(sl, Inches(2.3))

    # Use of funds
    tb(sl, MARGIN_L, Inches(2.6), Inches(4), Inches(0.3),
       "Use of funds", size=Pt(14), bold=True, color=BLACK)

    funds = [
        ("40%", "Engineering + mainnet deployment"),
        ("20%", "Security audit"),
        ("20%", "LP seed liquidity"),
        ("20%", "Mining pool partnerships"),
    ]
    for i, (pct, desc) in enumerate(funds):
        y = Inches(3.1) + i * Inches(0.55)
        tb(sl, MARGIN_L, y, Inches(0.8), Inches(0.3),
           pct, size=Pt(14), bold=True, color=ACCENT)
        tb(sl, Inches(1.8), y, Inches(4), Inches(0.3),
           desc, size=Pt(12), color=BODY)

    # What we deliver
    tb(sl, Inches(7.0), Inches(2.6), Inches(5), Inches(0.3),
       "This round delivers", size=Pt(14), bold=True, color=BLACK)

    deliverables = [
        "Security audit completed",
        "HashKey Chain mainnet deployment",
        "50 pilot miners onboarded",
        "$500K TVL",
    ]
    for i, d in enumerate(deliverables):
        y = Inches(3.1) + i * Inches(0.55)
        tb(sl, Inches(7.0), y, Inches(0.3), Inches(0.3),
           "→", size=Pt(12), bold=True, color=ACCENT)
        tb(sl, Inches(7.4), y, Inches(5), Inches(0.3),
           d, size=Pt(12), color=BODY)

    # From HashKey ecosystem
    h_rule(sl, Inches(5.5))
    tb(sl, MARGIN_L, Inches(5.8), Inches(5), Inches(0.3),
       "From HashKey ecosystem", size=Pt(14), bold=True, color=BLACK)

    asks = [
        "Technical partnership + infrastructure support",
        "HashKey Capital — investment + mining operator introductions",
        "HashKey Exchange — USDT liquidity integration + fiat on/off-ramp",
        "Compliance guidance — cross-jurisdictional pool withholding",
        "Post-hackathon incubation + ecosystem resource support",
    ]
    for i, a in enumerate(asks):
        y = Inches(6.1) + i * Inches(0.28)
        tb(sl, Inches(1.2), y, Inches(10), Inches(0.3),
           a, size=Pt(11), color=BODY)

    slide_num(sl, 14)


# ── 15. THANK YOU ────────────────────────────────────────────────────────────
def slide_thankyou(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    # Full-width accent bar
    add_rect(sl, Inches(0.45), Inches(0), Inches(0.12), SLIDE_H, fill_color=ACCENT)

    tb(sl, Inches(1.2), Inches(2.2), Inches(10), Inches(1.0),
       "Thank you.", size=Pt(52), bold=True, color=BLACK)

    h_rule(sl, Inches(3.5))

    links = [
        ("GitHub", "github.com/inchyangv/ctc-hashcredit"),
        ("Live Demo", "hashcredit.studioliq.com"),
        ("Contact", "inch@studioliq.com"),
    ]
    for i, (label, url) in enumerate(links):
        y = Inches(3.8) + i * Inches(0.5)
        tb(sl, Inches(1.2), y, Inches(2), Inches(0.3),
           label, size=Pt(13), bold=True, color=DARK)
        tb(sl, Inches(3.5), y, Inches(6), Inches(0.3),
           url, size=Pt(13), color=ACCENT)

    tb(sl, Inches(1.2), Inches(6.2), Inches(10), Inches(0.3),
       "HashKey Chain Testnet · chainId 133 · DeFi Track",
       size=Pt(11), color=MUTED)

    slide_num(sl, 15)


# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_cover(prs)         # 1
    slide_problem(prs)       # 2
    slide_market(prs)        # 3
    slide_insight(prs)       # 4
    slide_solution(prs)      # 5
    slide_how(prs)           # 6
    slide_biz_model(prs)     # 7
    slide_traction(prs)      # 8
    slide_why_now(prs)       # 9
    slide_competitive(prs)   # 10
    slide_roadmap(prs)       # 11
    slide_team(prs)          # 12
    slide_why_hashkey(prs)   # 13
    slide_ask(prs)           # 14
    slide_thankyou(prs)      # 15

    out = "HashCredit_PitchDeck.pptx"
    prs.save(out)
    print(f"Saved → {out}  ({TOTAL_SLIDES} slides)")


if __name__ == "__main__":
    main()
