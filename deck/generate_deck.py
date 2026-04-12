#!/usr/bin/env python3
"""
HashCredit Pitch Deck Generator — Neo-Brutalism Style
Generates a 16:9 PPTX with the design system from DECK.md.

Usage: python3 script/generate_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================
# Design Tokens
# ============================================

TEAL = RGBColor(0x0F, 0xA8, 0x9E)
BNB_YELLOW = RGBColor(0xF0, 0xB9, 0x0B)
SIGNAL_RED = RGBColor(0xFF, 0x3B, 0x30)
DEEP_BLACK = RGBColor(0x0A, 0x0F, 0x1C)
PURE_BLACK = RGBColor(0x00, 0x00, 0x00)
PURE_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF5, 0xF0)
WARM_GREY = RGBColor(0xF0, 0xED, 0xE8)
SOFT_WHITE = RGBColor(0xE0, 0xE8, 0xF0)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GREY = RGBColor(0x66, 0x66, 0x66)
LIGHT_TEAL = RGBColor(0xE8, 0xFA, 0xF8)

FONT_TITLE = "Arial Black"
FONT_BODY = "Courier New"
FONT_CAPTION = "Courier New"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BORDER_W = Pt(3)
SHADOW_OFF = Pt(6)


# ============================================
# Helpers
# ============================================

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, w, h, fill, border=PURE_BLACK, border_w=BORDER_W):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = border_w
    else:
        s.line.fill.background()
    return s


def add_shadow(slide, left, top, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + SHADOW_OFF, top + SHADOW_OFF, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = PURE_BLACK
    s.line.fill.background()
    return s


def card(slide, left, top, w, h, fill=PURE_WHITE):
    add_shadow(slide, left, top, w, h)
    return add_rect(slide, left, top, w, h, fill)


def tb(slide, left, top, w, h):
    return slide.shapes.add_textbox(left, top, w, h)


def set_tf(tf, text, font=FONT_BODY, size=Pt(14), color=NEAR_BLACK, bold=False, align=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align


def add_p(tf, text, font=FONT_BODY, size=Pt(14), color=NEAR_BLACK, bold=False, align=PP_ALIGN.LEFT, sp=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.name = font
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.space_before = sp
    return p


def accent_bar(slide, pos="top", color=TEAL):
    if pos == "top":
        r = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(10), color, border=None)
    elif pos == "left":
        r = add_rect(slide, Inches(0), Inches(0), Pt(10), SLIDE_H, color, border=None)
    elif pos == "bottom":
        r = add_rect(slide, Inches(0), SLIDE_H - Pt(10), SLIDE_W, Pt(10), color, border=None)


def wordmark(slide, color=PURE_WHITE):
    t = tb(slide, Inches(0.5), Inches(0.25), Inches(3), Inches(0.4))
    set_tf(t.text_frame, "HASHCREDIT", FONT_TITLE, Pt(11), color, bold=True)


def slide_num(slide, n, section, color=DARK_GREY):
    t = tb(slide, Inches(10.5), Inches(7.0), Inches(2.5), Inches(0.3))
    set_tf(t.text_frame, f"{n:02d} / {section.upper()}", FONT_CAPTION, Pt(9), color, align=PP_ALIGN.RIGHT)


def big_text(slide, text, left, top, color=NEAR_BLACK, size=Pt(96), align=PP_ALIGN.LEFT):
    t = tb(slide, left, top, Inches(12), Inches(1.8))
    set_tf(t.text_frame, text, FONT_TITLE, size, color, bold=True, align=align)
    return t


# ============================================
# Slides
# ============================================

def s01_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DEEP_BLACK)
    accent_bar(s, "top", TEAL)
    accent_bar(s, "bottom", TEAL)

    # Decorative teal block top-right
    add_rect(s, Inches(10.5), Inches(0.8), Inches(2.2), Inches(2.2), TEAL, border=None)

    big_text(s, "HASH", Inches(0.8), Inches(1.8), PURE_WHITE, Pt(120))
    big_text(s, "CREDIT", Inches(0.8), Inches(3.2), TEAL, Pt(120))

    t = tb(s, Inches(0.8), Inches(5.2), Inches(9), Inches(0.6))
    set_tf(t.text_frame, "Revenue-based financing for Bitcoin miners. Pool-enforced. SPV-proven.", FONT_BODY, Pt(18), SOFT_WHITE)

    t = tb(s, Inches(0.8), Inches(6.2), Inches(6), Inches(0.4))
    set_tf(t.text_frame, "Built on BNB Chain  |  bnb.hash.credit", FONT_BODY, Pt(14), BNB_YELLOW)

    slide_num(s, 1, "COVER", DARK_GREY)


def s02_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, PURE_WHITE)
    accent_bar(s, "left", SIGNAL_RED)
    wordmark(s, NEAR_BLACK)

    big_text(s, "$17B", Inches(0.8), Inches(0.6), SIGNAL_RED, Pt(100))
    t = tb(s, Inches(0.8), Inches(2.3), Inches(11), Inches(0.5))
    set_tf(t.text_frame, "IN ANNUAL MINING REVENUE. ZERO ON-CHAIN CREDIT.", FONT_TITLE, Pt(22), NEAR_BLACK, bold=True)

    data = [
        ("$137K", "Full cost per BTC\npost-halving. ROI > 1,200 days."),
        ("$11B+", "Mining debt raised\nsince 2023. Demand is proven."),
        ("BANKRUPT", "BlockFi, Celsius, Genesis.\nTrust model failed."),
    ]
    for i, (num, desc) in enumerate(data):
        left = Inches(0.8 + i * 4.1)
        c = card(s, left, Inches(3.2), Inches(3.6), Inches(3.4))
        tf = c.text_frame
        tf.word_wrap = True
        set_tf(tf, num, FONT_TITLE, Pt(40), SIGNAL_RED, bold=True)
        add_p(tf, desc, FONT_BODY, Pt(14), NEAR_BLACK, sp=Pt(16))
        # Teal underline inside card
        add_rect(s, left + Pt(20), Inches(4.3), Inches(2.8), Pt(4), TEAL, border=None)

    slide_num(s, 2, "PROBLEM")


def s03_insight(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, TEAL)
    accent_bar(s, "top", DEEP_BLACK)
    wordmark(s, PURE_WHITE)

    # Big statement — two lines
    big_text(s, "MATH,", Inches(0.8), Inches(1.2), PURE_WHITE, Pt(96))
    big_text(s, "NOT TRUST.", Inches(0.8), Inches(2.8), DEEP_BLACK, Pt(96))

    # Explanation in a dark card
    c = card(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.2), DEEP_BLACK)
    tf = c.text_frame
    tf.word_wrap = True
    set_tf(tf, "Pool payouts = Bitcoin txs committed to by PoW. Payout history IS the hashrate record.", FONT_BODY, Pt(16), PURE_WHITE)
    add_p(tf, "SPV turns it into trustless on-chain evidence. No oracle. No custodian.", FONT_BODY, Pt(16), TEAL, sp=Pt(12))
    add_p(tf, "Stripe Capital underwrites with payment data. We use cryptographic proof of mining revenue.", FONT_BODY, Pt(14), SOFT_WHITE, sp=Pt(12))

    slide_num(s, 3, "INSIGHT", PURE_WHITE)


def s04_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DEEP_BLACK)
    accent_bar(s, "top", TEAL)
    wordmark(s, PURE_WHITE)

    t = tb(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
    set_tf(t.text_frame, "STRIPE CAPITAL FOR BITCOIN MINERS.", FONT_TITLE, Pt(36), PURE_WHITE, bold=True)

    steps = ["Pool payout\n(BTC tx)", "SPV proof\n(trustless)", "Credit limit\nupdates", "Draw USDT\n(BNB Chain)", "Pool auto-\nwithholds"]
    for i, step in enumerate(steps):
        left = Inches(0.4 + i * 2.55)
        c = card(s, left, Inches(2.0), Inches(2.2), Inches(2.0), TEAL)
        tf = c.text_frame
        tf.word_wrap = True
        set_tf(tf, f"0{i+1}", FONT_TITLE, Pt(32), PURE_WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_p(tf, step, FONT_BODY, Pt(12), PURE_WHITE, align=PP_ALIGN.CENTER, sp=Pt(8))

    # Arrow line
    add_rect(s, Inches(0.8), Inches(4.3), Inches(11.5), Pt(3), TEAL, border=None)

    # Bottom bullets
    bullets = [
        "No BTC lockup — mine through your registered pool",
        "Default → pool redirects miner's hashrate. No courts.",
        "Modular IVerifierAdapter — swap proof sources, credit logic untouched",
    ]
    for i, b in enumerate(bullets):
        t = tb(s, Inches(0.8), Inches(4.8 + i * 0.65), Inches(11), Inches(0.5))
        set_tf(t.text_frame, f"→  {b}", FONT_BODY, Pt(14), SOFT_WHITE)

    slide_num(s, 4, "SOLUTION", DARK_GREY)


def s05_how(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    accent_bar(s, "left", TEAL)
    wordmark(s, NEAR_BLACK)

    t = tb(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_tf(t.text_frame, "HOW IT WORKS", FONT_TITLE, Pt(40), NEAR_BLACK, bold=True)

    steps = [
        ("01", "REGISTER", "Pool agrees to withhold repayment"),
        ("02", "MINE", "Miner gets paid by pool (BTC tx)"),
        ("03", "DETECT", "Off-chain worker detects payout"),
        ("04", "PROVE", "Build SPV proof — headers + Merkle"),
        ("05", "VERIFY", "On-chain PoW + Merkle + output check"),
        ("06", "CREDIT", "Credit limit auto-updates"),
        ("07", "DRAW", "Miner draws USDT / Pool withholds"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(1.3 + i * 0.82)
        # Number block
        add_rect(s, Inches(0.8), y, Inches(0.9), Inches(0.65), TEAL if i < 5 else DEEP_BLACK, border=None)
        t = tb(s, Inches(0.8), y + Pt(2), Inches(0.9), Inches(0.55))
        set_tf(t.text_frame, num, FONT_TITLE, Pt(20), PURE_WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Title
        t = tb(s, Inches(2.0), y + Pt(2), Inches(2.5), Inches(0.55))
        set_tf(t.text_frame, title, FONT_TITLE, Pt(16), NEAR_BLACK, bold=True)
        # Desc
        t = tb(s, Inches(4.8), y + Pt(4), Inches(8), Inches(0.5))
        set_tf(t.text_frame, desc, FONT_BODY, Pt(13), DARK_GREY)

    slide_num(s, 5, "MECHANISM")


def s06_why_now(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, PURE_WHITE)
    accent_bar(s, "left", SIGNAL_RED)
    wordmark(s, NEAR_BLACK)

    t = tb(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_tf(t.text_frame, "WHY NOW", FONT_TITLE, Pt(44), NEAR_BLACK, bold=True)

    signals = [
        ("HALVING", "Revenue per hash cut 50%.\nWorking capital is existential."),
        ("$11B DEBT", "Miners raised $11B+ since 2023.\nDemand is proven & structural."),
        ("1 ZH/s", "Hardware ROI > 1,200 days.\nMining is capital-intensive."),
        ("TRUST FAILED", "BlockFi, Celsius, Genesis —\nall bankrupt. Trustless is next."),
        ("RWA $17B+", "90% is static (treasuries).\nProductive RWA is the gap."),
        ("UNOCCUPIED", "No one does SPV proof +\npool enforcement + on-chain credit."),
    ]
    for i, (kw, desc) in enumerate(signals):
        row, col = i // 3, i % 3
        left = Inches(0.8 + col * 4.1)
        top = Inches(1.5 + row * 2.8)
        c = card(s, left, top, Inches(3.6), Inches(2.3))
        tf = c.text_frame
        tf.word_wrap = True
        set_tf(tf, kw, FONT_TITLE, Pt(24), SIGNAL_RED, bold=True)
        add_p(tf, desc, FONT_BODY, Pt(13), NEAR_BLACK, sp=Pt(12))

    slide_num(s, 6, "TIMING")


def s07_market(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    accent_bar(s, "left", TEAL)
    wordmark(s, NEAR_BLACK)

    big_text(s, "$17.2B", Inches(0.8), Inches(0.3), NEAR_BLACK, Pt(80))
    t = tb(s, Inches(0.8), Inches(1.8), Inches(8), Inches(0.4))
    set_tf(t.text_frame, "ANNUAL BTC MINER REVENUE (2025, THE BLOCK)", FONT_TITLE, Pt(14), DARK_GREY)

    # Three cards
    cards_data = [
        ("SAM", "$5.2–6.9B", "Mid-market miners\n30-40% of hashrate\nUnderserved by all options"),
        ("DISTRIBUTION", "TOP 10 POOLS", "= 90%+ of hashrate\n1 pool = 1000s borrowers\nB2B, not B2C"),
        ("COMPARABLE", "STRIPE CAPITAL", "$9B+ merchant advances\nSame RBF model\nMining = stronger data"),
    ]
    for i, (label, big, desc) in enumerate(cards_data):
        left = Inches(0.8 + i * 4.1)
        c = card(s, left, Inches(2.5), Inches(3.6), Inches(3.0))
        tf = c.text_frame
        tf.word_wrap = True
        set_tf(tf, label, FONT_CAPTION, Pt(10), DARK_GREY)
        add_p(tf, big, FONT_TITLE, Pt(26), TEAL, bold=True, sp=Pt(6))
        add_p(tf, desc, FONT_BODY, Pt(13), NEAR_BLACK, sp=Pt(12))

    # LP yield bar
    c = card(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0), DEEP_BLACK)
    tf = c.text_frame
    tf.word_wrap = True
    set_tf(tf, "LP: 8% APR  |  PancakeSwap 2-4%  |  Venus 3-6%  |  Real yield, not emissions", FONT_BODY, Pt(15), BNB_YELLOW, bold=True, align=PP_ALIGN.CENTER)

    slide_num(s, 7, "MARKET")


def s08_competitive(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, PURE_WHITE)
    accent_bar(s, "left", TEAL)
    wordmark(s, NEAR_BLACK)

    t = tb(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_tf(t.text_frame, "ONLY TRUSTLESS + POOL-ENFORCED", FONT_TITLE, Pt(32), NEAR_BLACK, bold=True)

    headers = ["", "HASHCREDIT", "EVERYONE ELSE"]
    rows = [
        ["COLLATERAL", "None — revenue-proven", "BTC, ASICs, financials"],
        ["VERIFICATION", "SPV — pure math", "Oracle / manual review"],
        ["REPAYMENT", "Auto — pool withholds", "Manual, trust-dependent"],
        ["SPEED", "Instant, auto-updating", "Days to weeks"],
        ["ACCESS", "Permissionless via pool", "KYC, minimums, gates"],
    ]

    col_w = [Inches(2.8), Inches(4.5), Inches(4.5)]
    col_x = [Inches(0.8), Inches(3.7), Inches(8.3)]

    # Header
    for j, h in enumerate(headers):
        bg = [DEEP_BLACK, TEAL, NEAR_BLACK][j]
        r = add_rect(s, col_x[j], Inches(1.4), col_w[j], Inches(0.7), bg)
        tf = r.text_frame
        set_tf(tf, h, FONT_TITLE, Pt(14), PURE_WHITE, bold=True, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            top = Inches(2.2 + i * 0.95)
            bg = PURE_WHITE if j != 1 else LIGHT_TEAL
            r = add_rect(s, col_x[j], top, col_w[j], Inches(0.8), bg, border_w=Pt(1.5))
            tf = r.text_frame
            tf.word_wrap = True
            color = TEAL if j == 1 else NEAR_BLACK
            set_tf(tf, cell, FONT_BODY, Pt(12), color, bold=(j == 0), align=PP_ALIGN.CENTER)

    slide_num(s, 8, "EDGE")


def s09_traction(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    accent_bar(s, "left", TEAL)
    wordmark(s, NEAR_BLACK)

    # Big "7" with teal block behind
    add_rect(s, Inches(0.8), Inches(0.5), Inches(2.2), Inches(2.2), TEAL, border=None)
    t = tb(s, Inches(0.8), Inches(0.5), Inches(2.2), Inches(2.2))
    set_tf(t.text_frame, "7", FONT_TITLE, Pt(120), PURE_WHITE, bold=True, align=PP_ALIGN.CENTER)

    t = tb(s, Inches(3.4), Inches(1.0), Inches(9), Inches(0.8))
    set_tf(t.text_frame, "CONTRACTS LIVE\nON BSC TESTNET", FONT_TITLE, Pt(32), NEAR_BLACK, bold=True)

    items = [
        "SPV proofs from real Bitcoin testnet transactions",
        "Full borrow / repay lifecycle operational",
        "24/7 automated prover worker running",
        "Frontend dashboard live at bnb.hash.credit",
        "Unit + integration + invariant fuzzing + gas profiling",
        "Modular IVerifierAdapter — plug-and-prove architecture",
    ]
    for i, item in enumerate(items):
        row, col = i // 2, i % 2
        left = Inches(0.8 + col * 6.2)
        top = Inches(3.2 + row * 1.3)
        c = card(s, left, top, Inches(5.6), Inches(0.95))
        tf = c.text_frame
        tf.word_wrap = True
        set_tf(tf, f"✓  {item}", FONT_BODY, Pt(13), NEAR_BLACK)

    slide_num(s, 9, "TRACTION")


def s10_rwa(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, TEAL)
    accent_bar(s, "top", DEEP_BLACK)
    wordmark(s, PURE_WHITE)

    t = tb(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7))
    set_tf(t.text_frame, "THE MISSING PIECE IN RWA", FONT_TITLE, Pt(40), PURE_WHITE, bold=True)

    # Left: Existing RWA
    c = card(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8), PURE_WHITE)
    tf = c.text_frame
    tf.word_wrap = True
    set_tf(tf, "EXISTING RWA", FONT_TITLE, Pt(22), SIGNAL_RED, bold=True)
    add_p(tf, "Ondo / Maple / Centrifuge", FONT_BODY, Pt(11), DARK_GREY, sp=Pt(4))
    for item in ["Tokenized treasuries — static yield", "Oracle + legal attestation required", "Quarterly / annual revenue cycles", "Jurisdiction-dependent enforcement"]:
        add_p(tf, f"✗  {item}", FONT_BODY, Pt(14), NEAR_BLACK, sp=Pt(14))

    # Right: HashCredit
    c = card(s, Inches(7), Inches(1.6), Inches(5.5), Inches(4.8), DEEP_BLACK)
    tf = c.text_frame
    tf.word_wrap = True
    set_tf(tf, "HASHCREDIT", FONT_TITLE, Pt(22), TEAL, bold=True)
    add_p(tf, "Bitcoin mining revenue", FONT_BODY, Pt(11), DARK_GREY, sp=Pt(4))
    for item in ["Mining revenue — dynamic daily yield", "SPV proof — pure cryptography", "Daily / weekly payout cycles", "Pool enforcement — global, no courts"]:
        add_p(tf, f"✓  {item}", FONT_BODY, Pt(14), SOFT_WHITE, sp=Pt(14))

    # Bottom tagline
    t = tb(s, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
    set_tf(t.text_frame, "Accounts receivable factoring for hashrate. Stripe Capital model + cryptographic trust.", FONT_BODY, Pt(14), PURE_WHITE, align=PP_ALIGN.CENTER)

    slide_num(s, 10, "RWA", PURE_WHITE)


def s11_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, OFF_WHITE)
    accent_bar(s, "left", TEAL)
    wordmark(s, NEAR_BLACK)

    t = tb(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_tf(t.text_frame, "ROADMAP", FONT_TITLE, Pt(44), NEAR_BLACK, bold=True)

    phases = [
        ("Q1 2026", "DONE", "7 contracts live\nSPV verifier\nFull lifecycle\nFrontend on BSC Testnet", TEAL, PURE_WHITE),
        ("Q2 2026", "NOW", "Security audit\nBSC mainnet deploy\n10 pilot miners", BNB_YELLOW, NEAR_BLACK),
        ("Q3 2026", "NEXT", "50 miners\n$500K TVL\nPool API partnerships", PURE_WHITE, NEAR_BLACK),
        ("Q4 2026", "VISION", "$5M TVL\n200+ miners\nSeries A", PURE_WHITE, NEAR_BLACK),
    ]
    for i, (when, status, desc, bg, tc) in enumerate(phases):
        left = Inches(0.5 + i * 3.15)
        c = card(s, left, Inches(1.6), Inches(2.8), Inches(5.2), bg)
        tf = c.text_frame
        tf.word_wrap = True
        set_tf(tf, when, FONT_TITLE, Pt(22), tc, bold=True)
        # Status badge
        add_p(tf, status, FONT_TITLE, Pt(14), tc, bold=True, sp=Pt(6))
        add_p(tf, desc, FONT_BODY, Pt(13), tc, sp=Pt(20))

    slide_num(s, 11, "ROADMAP")


def s12_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WARM_GREY)
    accent_bar(s, "top", TEAL)
    wordmark(s, NEAR_BLACK)

    t = tb(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
    set_tf(t.text_frame, "TEAM", FONT_TITLE, Pt(44), NEAR_BLACK, bold=True)

    # CEO
    c = card(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0))
    tf = c.text_frame
    tf.word_wrap = True
    set_tf(tf, "INCHEOL YANG", FONT_TITLE, Pt(26), NEAR_BLACK, bold=True)
    add_p(tf, "CEO  /  KAIST CS", FONT_TITLE, Pt(14), TEAL, bold=True, sp=Pt(4))
    add_p(tf, "", size=Pt(6), sp=Pt(2))
    for line in [
        "Co-founded DeFi system trading house",
        "$20M managed at 40%+ APR",
        "",
        "KRAFTON PUBG Studio",
        "  In-game trading & payment systems",
        "",
        "Coinone Exchange (#3 Korea)",
        "  Smart order routing, institutional API",
    ]:
        add_p(tf, line, FONT_BODY, Pt(12), NEAR_BLACK if line and not line.startswith(" ") else DARK_GREY, sp=Pt(3))

    # CTO
    c = card(s, Inches(7), Inches(1.6), Inches(5.5), Inches(5.0))
    tf = c.text_frame
    tf.word_wrap = True
    set_tf(tf, "JUHYEONG PARK", FONT_TITLE, Pt(26), NEAR_BLACK, bold=True)
    add_p(tf, "CTO  /  Yonsei CS", FONT_TITLE, Pt(14), TEAL, bold=True, sp=Pt(4))
    add_p(tf, "", size=Pt(6), sp=Pt(2))
    for line in [
        "CTO of Onther",
        "  Led mainnet to $750M market cap",
        "  Designed Plasma EVM",
        "",
        "Chainpartners",
        "  DEX aggregator design",
        "  Perpetual DEX architecture",
    ]:
        add_p(tf, line, FONT_BODY, Pt(12), NEAR_BLACK if line and not line.startswith(" ") else DARK_GREY, sp=Pt(3))

    t = tb(s, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4))
    set_tf(t.text_frame, "Two founders built the entire protocol — contracts, SPV prover, worker, frontend.", FONT_BODY, Pt(13), DARK_GREY, align=PP_ALIGN.CENTER)

    slide_num(s, 12, "TEAM")


def s13_why_bnb(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DEEP_BLACK)
    wordmark(s, PURE_WHITE)

    # Left panel — BNB Yellow
    add_rect(s, Inches(0), Inches(0), Inches(5.3), SLIDE_H, BNB_YELLOW, border=None)
    # Divider
    add_rect(s, Inches(5.3), Inches(0), Pt(4), SLIDE_H, PURE_WHITE, border=None)

    # Left: What we bring
    t = tb(s, Inches(0.5), Inches(0.5), Inches(4.5), Inches(0.6))
    set_tf(t.text_frame, "WHAT WE BRING", FONT_TITLE, Pt(28), NEAR_BLACK, bold=True)

    left_items = [
        "Real miner TVL — not synthetic",
        "Novel RWA: productive asset revenue",
        "Organic USDT demand ($17B sector)",
        "Predictable daily on-chain traffic",
    ]
    for i, item in enumerate(left_items):
        t = tb(s, Inches(0.5), Inches(1.4 + i * 0.7), Inches(4.5), Inches(0.5))
        set_tf(t.text_frame, f"→  {item}", FONT_BODY, Pt(13), NEAR_BLACK)

    # Right: Why BNB
    t = tb(s, Inches(5.8), Inches(0.5), Inches(7), Inches(0.6))
    set_tf(t.text_frame, "WHY BNB CHAIN", FONT_TITLE, Pt(28), BNB_YELLOW, bold=True)

    right_items = [
        "Deepest USDT liquidity in DeFi",
        "~$0.01 gas (100x < ETH L1)",
        "Largest RWA ecosystem incentives",
        "Global user base — Asia mining hub",
    ]
    for i, item in enumerate(right_items):
        t = tb(s, Inches(5.8), Inches(1.4 + i * 0.7), Inches(7), Inches(0.5))
        set_tf(t.text_frame, f"→  {item}", FONT_BODY, Pt(13), SOFT_WHITE)

    # Nano Labs callout card
    c = card(s, Inches(5.8), Inches(4.2), Inches(6.8), Inches(2.8), DEEP_BLACK)
    c.line.color.rgb = BNB_YELLOW
    tf = c.text_frame
    tf.word_wrap = True
    set_tf(tf, "NANO LABS × HASHCREDIT", FONT_TITLE, Pt(20), BNB_YELLOW, bold=True)
    add_p(tf, "Nano Labs = world's largest mining chip designer", FONT_BODY, Pt(13), SOFT_WHITE, sp=Pt(10))
    add_p(tf, "Their chips → ASICs → miners → pools → HashCredit → BNB Chain", FONT_BODY, Pt(13), TEAL, sp=Pt(8))
    add_p(tf, "Vertical integration: every actor benefits.", FONT_BODY, Pt(14), PURE_WHITE, bold=True, sp=Pt(12))

    # Nano Labs on left panel
    c2 = card(s, Inches(0.5), Inches(4.5), Inches(4.3), Inches(2.5), NEAR_BLACK)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    set_tf(tf2, "VALUE CHAIN", FONT_TITLE, Pt(16), BNB_YELLOW, bold=True)
    add_p(tf2, "Chip Designer (Nano Labs)\n  ↓\nMiner\n  ↓\nPool\n  ↓\nHashCredit\n  ↓\nBNB Chain", FONT_BODY, Pt(13), PURE_WHITE, sp=Pt(8))

    slide_num(s, 13, "BNB CHAIN", DARK_GREY)


def s14_ask(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, DEEP_BLACK)
    accent_bar(s, "top", TEAL)
    accent_bar(s, "bottom", TEAL)
    wordmark(s, PURE_WHITE)

    big_text(s, "$250K SEED", Inches(0.8), Inches(0.5), PURE_WHITE, Pt(72))

    funds = [
        ("40%", "$100K", "Engineering +\nMainnet Deploy"),
        ("20%", "$50K", "Security\nAudit"),
        ("20%", "$50K", "Vault USDT\nLiquidity"),
        ("20%", "$50K", "Pool Partners\n+ GTM"),
    ]
    for i, (pct, amt, desc) in enumerate(funds):
        left = Inches(0.5 + i * 3.15)
        c = card(s, left, Inches(2.0), Inches(2.8), Inches(2.4), TEAL)
        tf = c.text_frame
        tf.word_wrap = True
        set_tf(tf, pct, FONT_TITLE, Pt(40), PURE_WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_p(tf, amt, FONT_TITLE, Pt(16), BNB_YELLOW, bold=True, align=PP_ALIGN.CENTER, sp=Pt(4))
        add_p(tf, desc, FONT_BODY, Pt(12), PURE_WHITE, align=PP_ALIGN.CENTER, sp=Pt(8))

    # Event incentives
    t = tb(s, Inches(0.8), Inches(4.7), Inches(11), Inches(0.4))
    set_tf(t.text_frame, "FROM RWA DEMO DAY:", FONT_TITLE, Pt(18), BNB_YELLOW, bold=True)

    incentives = [
        ("ICC Incubation", "$100K value — acceleration & advisory"),
        ("BNB Chain RWA Fast-Track", "Liquidity seeding, TVL incentives, tech guidance"),
        ("Nano Labs Mentorship", "World's largest mining chip designer — direct advisory"),
        ("HK Web3 Festival", "Showcase slot — investor & partner exposure"),
    ]
    for i, (title, desc) in enumerate(incentives):
        t = tb(s, Inches(0.8), Inches(5.2 + i * 0.45), Inches(11), Inches(0.4))
        set_tf(t.text_frame, f"{title}  —  {desc}", FONT_BODY, Pt(12), SOFT_WHITE)

    slide_num(s, 14, "THE ASK", DARK_GREY)


# ============================================
# Main
# ============================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s01_cover(prs)
    s02_problem(prs)
    s03_insight(prs)
    s04_solution(prs)
    s05_how(prs)
    s06_why_now(prs)
    s07_market(prs)
    s08_competitive(prs)
    s09_traction(prs)
    s10_rwa(prs)
    s11_roadmap(prs)
    s12_team(prs)
    s13_why_bnb(prs)
    s14_ask(prs)

    out = "HashCredit_Deck_RWA_DemoDay.pptx"
    prs.save(out)
    print(f"Generated: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
